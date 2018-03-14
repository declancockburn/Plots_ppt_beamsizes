# -*- coding: utf-8 -*-
"""
Created on Mon Feb 12 11:57:13 2018
Python 3.6

@author: dcockbur
"""
#%% 
#Preamble

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import os
from pptx import Presentation
from pptx.util import Cm #centimetres
plt.style.use('seaborn') #use seaborn standard style
plt.ioff() #turn off display plots unless using plt.show()

#%% VARIABLES

path = 'C:\\Users\\dcockbur\\Documents\\Python\\Scripts For Work\\Beamsizeplots'
datasubfolder = '\\data\\'


#%% Make a lis tof filenames:

filelist = [f for f in os.listdir(path + datasubfolder) if f.endswith('.csv')] # make a list of files
os.chdir(path) #set the working directory to path

## Potential alternative to use to select all files
#import glob
#path = r'C:\Users\dcockbur\Documents\Python\Scripts For Work\Beamsizeplots'
#allfiles = glob.glob(path + '/*.csv')
#print(allfiles)

#%% put data into MP and PP dfs and CSVs.
labels = [x[:-4] for x in filelist] #remove the .csv from the filelist names
labeldic = dict(zip(filelist,labels)) #make a dictionary of labels taking away the .csv

temp = pd.read_csv('data\\'+filelist[0]) #take an example file, first one, as template.

MPdf = pd.DataFrame()#blank dataframe
PPdf = pd.DataFrame()
MPdf['Location'] = temp['Location'] #put in the first PA0-in etc. column
PPdf['Location'] = temp['Location']

for file_ in filelist: #populate the columns
    frame = pd.read_csv(path + '\\data\\' + file_, index_col=None) #make a temp df
    MPdf[file_] = frame['MP'] # Put the MP data into the MP only df, names by the machinename
    PPdf[file_] = frame['PP']
MPdf.rename(columns = labeldic, inplace = True) # Rename columns without the csv name, with the dictionary
PPdf.rename(columns = labeldic, inplace = True)


MPdf.to_csv('MP.csv') #Make a new CSV file just for the the MP and PP data
PPdf.to_csv('PP.csv')

#%%
#Set the Location (PA) as the index and transpose the DF to make the location the column names

MPdf.set_index('Location', inplace=True)
PPdf.set_index('Location', inplace=True)
MPdfT = MPdf.T
PPdfT = PPdf.T

#%%
#Get statistics on the transposed data, specifically the 2 sigma, 3 sigma (according to wiki)

MPstats = MPdfT.describe(percentiles = [0.0013, 0.0228, 0.5, 0.9772, 0.9987]).T
PPstats = PPdfT.describe(percentiles = [0.0013, 0.0228, 0.5, 0.9772, 0.9987]).T
MPstats.reset_index(inplace=True)
PPstats.reset_index(inplace=True)


#%% Reset the index and remove the first "index" column.
MPdf.reset_index(inplace=True)
MPdf = MPdf.iloc[:,1:]
PPdf.reset_index(inplace=True)
PPdf = PPdf.iloc[:,1:]

#%%
# Funtion to calculate the "slope" of the beamsize convergance/divergance
def slopes(beamdf):
    slope = beamdf.iloc[-1] - beamdf.iloc[-2]
    slopemean = float(np.mean(slope))
    slopeID = slope.astype('str')
    for i in range(len(slope)):
        if slope[i] < slopemean/0.9:
            slopeID[i] = 'Converging'
        elif slope[i] > slopemean*0.9:
            slopeID[i] = 'Diverging'
        else:
            slopeID[i] = ''            
    #slopeID = ['Converging' for x > slopemean in slope else: 'Diverging']
    #return slope,slopemean
    return slope, slopemean, slopeID
slopes(MPdf)
#slope, slopemean = slopes(MPdf)
#print(slope)
#MPslope.head(25)
#MPslopemean
#MPslope['DL24_PR7']

#%%

close('all') #Close all figs in case any are open

#colours RGB values
red = (0.99,0.01,0.01) 
orange = (0.99,0.4,0.0)
green = (0.1,0.8,0.0)

def figback(df): #function to make the figure's "background", with the 2sigma boundaries and colours. arg is stats df
    fig1 = plt.figure(figsize=(6,4.8)) #define figure
    ax1 = fig1.add_subplot(111) #define axis on 
    X = df.index.values # Makes a numerical X axis, of the same length
    ax1.plot(df['mean'], 'g--', linewidth=1) # Plot the mean
    
    ax1.plot(df['2.3%'], '--', c=red, linewidth=1, label = "2"+r"$\sigma$"+" (population)") #Plot the lower boundary
    ax1.plot(df['97.7%'], '--', c=red, linewidth=1, label='') #Plot the upper boundary
    
    ax1.fill_between(X,df['2.3%'],df['97.7%'], color=green, alpha = 0.15, linewidth=0.4) #fill center between two boundarys
    ax1.fill_between(X,0,df['2.3%'], color=red, alpha = 0.15, linewidth=0) #fill 0 to lower boundary
    ax1.fill_between(X,df['97.7%'], 40, color=red, alpha = 0.15, linewidth=0) # fill upper boundary to 40 
    
    plt.ylim(np.min(df['min']),np.max(df['max'])) #set the y-range
    plt.xlim(X[0],X[-1]) #set the x range as the same as the data for aesthetics.
    ax1.set_xticks(X) #put ticks where the X-values are
    ax1.set_xticklabels(df['Location'], rotation = 30, ha = 'right') #Label the ticks, rotate and right align them
#    ax1.legend()
#    fig1.tight_layout()
    return fig1


#%%
machines = list(MPdf.columns.values)
#for machine in machines:
#    ax1.plot(MPdf[machine],'--', linewidth = 1)

close('all')
#Function to make the figures and save them
def makefigs(df, beamdf, name): #the input args are the statistics dataframe, the MP or PP dataframe, and the name MP or PP as a string.
    slope,slopemean,slopeID = slopes(beamdf) #call function that calculates slopes
    for machine in machines: #For each machine
        fig2 = figback(df) #call the background figure
        ax1 = fig2.add_subplot(111) #Declare the axis
        #machine = machines[0]
        ax1.plot(beamdf[machine],'-', linewidth = 1.5) # plot the beamsizes on the figure
        
        if name == 'MP': #to add a label if converging or diverging, only for MP
            kleurdic = {'Converging':'green', 'Diverging':'red','':'red'} #dictionary of colours depending on slopeID
            ax1.text(6.8,16, '{}'.format(slopeID[machine]), horizontalalignment='left', 
                     bbox = {'facecolor': kleurdic[slopeID[machine]], 'alpha':0.5, 'pad':5}) #add a text box with converging or diverging (or no) label
        ax1.text(6.8,14,'PA3 Slope = {:.2f}\nAvg Slope = {:.2f}'.format(slope[machine],slopemean), horizontalalignment='left') #add a text box with slope and avg slope
        
        ax1.legend() # put in a legend 
        plt.title('{} {} ({})'.format(name, machine.split('_')[0], machine.split('_')[1]), fontsize = 16)
        fig2.tight_layout() # fit the layout better
        
        #fig2.show() #If you want to show the figures.
        fig2.savefig('{}_{}.png'.format(machine, name)) #save each figure
        plt.close(fig2) #Properly closes figure so it does not accumulate in memory

makefigs(MPstats, MPdf, 'MP') # Run the function above for the stats, the data and the label/name.
makefigs(PPstats, PPdf, 'PP')


#%% Putting it all together in a power point presentation

import datetime 
now = datetime.datetime.now() # Get the current date for the title slide


# Using template slide, make title slide and save
pptpath = 'template.pptx'
prs = Presentation(pptpath)

# Make the title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title_name = slide.shapes.placeholders[10]
title_loc = slide.shapes.placeholders[11]
title = slide.shapes.placeholders[12]
title_sub = slide.shapes.placeholders[14]
title_name.text = 'Declan Cockburn'
title_loc.text = '{}/{}/{}, Hsinchu Taiwan'.format(now.year, now.month, now.day)
title_sub.text = ' '
title.text = 'Alignment Beamsizes in S3 PAs'
prs.save('S3_PA_Alignment_beamsizes.pptx')


# Using content slide make first content slide with text.
prs = Presentation('S3_PA_Alignment_beamsizes.pptx')
insert_slide = prs.slide_layouts[1]
slide = prs.slides.add_slide(insert_slide)
shapes = slide.shapes
title2 = shapes.title
subtitle2 = shapes.placeholders[14]
text2 = shapes.placeholders[13]
title2.text = 'Introduction'
subtitle2.text = 'and Motivation'
text2.text = """
Below are the PA alignment beamsize plots from S3 SAT, with 2-Sigma of the population as boundaries (2.28 - 97.72 percentile) to show outliers. For MP, considered more relevant, an arbitrary >10% deviation from the mean slope (between PA3-in and PA3-out) is chosen here to indicate whether a beam will be "converging" or "diverging" from the mean.

Until another metric is available (PA3 wavefront) this:
- Can be used as part of the acceptance criteria for SATs (and feedback for their toroid config).
- May help predict the choice of M30 beam expander, saving time.
- May help predict expected beam-size on M130.
"""
prs.save('S3_PA_Alignment_beamsizes.pptx')

# Populate the rest of the content slides with data/pictures
for machine in machines:
    prs = Presentation('S3_PA_Alignment_beamsizes.pptx')
    blank_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes
    title2 = shapes.title
    subtitle2 = shapes.placeholders[14]
    text2 = shapes.placeholders[13]
    subtitle2.text = ' '
    text2.text = ' '
    title2.text = '{} ({})'.format(machine.split('_')[0], machine.split('_')[1])
    MP_img_path = '{}_MP.png'.format(machine)
    PP_img_path = '{}_PP.png'.format(machine)
    
    #Image 1 placement
    left = Cm(1)
    top = Cm(3)
    height = Cm(10)
    pic = slide.shapes.add_picture(MP_img_path, left, top, height = height)
    
    #Image 2 placement
    left = Cm(13.5)
    top = Cm(4.8)
    height = Cm(8)
    pic = slide.shapes.add_picture(PP_img_path, left, top, height=height)
    
    prs.save('S3_PA_Alignment_beamsizes.pptx')
    
#%% Uncomment to see the placeholders available in the template slide
#    
#prs = Presentation('template.pptx')
#slide = prs.slides.add_slide(prs.slide_layouts[1])
##slide.placeholder[0].name
#for shape in slide.placeholders:
#    print('%d ---- %s' % (shape.placeholder_format.idx, shape.name))


#%%
